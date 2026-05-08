import os
from database.engine import get_connection


def extract_text(file_path: str, file_type: str) -> str:
    if file_type == 'pdf':
        return _extract_pdf(file_path)
    elif file_type in ('docx', 'doc'):
        return _extract_docx(file_path)
    elif file_type in ('pptx', 'ppt'):
        return _extract_pptx(file_path)
    elif file_type in ('txt', 'md'):
        return _extract_txt(file_path)
    else:
        return ''


def _extract_pdf(file_path: str) -> str:
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return '\n'.join(text_parts)
    except Exception as e:
        return ''


def _extract_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ''


def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
        return '\n'.join(text_parts)
    except Exception:
        return ''


def _extract_txt(file_path: str) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()
        except Exception:
            return ''


def process_text_extraction(item_id: str):
    conn = get_connection()
    try:
        row = conn.execute("SELECT * FROM items WHERE id = ?", [item_id]).fetchone()
        if not row:
            return

        full_path = _resolve_path(row['file_path'])
        if not os.path.exists(full_path):
            conn.execute("UPDATE items SET status = 'error', error_msg = '文件不存在' WHERE id = ?", [item_id])
            conn.commit()
            return

        conn.execute("UPDATE items SET status = 'processing' WHERE id = ?", [item_id])

        text = extract_text(full_path, row['file_type'])

        conn.execute(
            "UPDATE items SET extracted_text = ?, status = 'done', updated_at = datetime('now') WHERE id = ?",
            [text, item_id]
        )

        conn.execute(
            "UPDATE processing_jobs SET status = 'done', finished_at = datetime('now') WHERE item_id = ? AND task_type = 'text_extract'",
            [item_id]
        )
        conn.commit()
    except Exception as e:
        conn.execute(
            "UPDATE items SET status = 'error', error_msg = ? WHERE id = ?",
            [str(e), item_id]
        )
        conn.commit()
    finally:
        conn.close()


def _resolve_path(relative_path: str) -> str:
    import config
    return os.path.join(config.FILES_DIR, relative_path)
